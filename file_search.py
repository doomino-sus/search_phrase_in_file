import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import datetime
from docx import Document
import pandas as pd
from pptx import Presentation
import PyPDF2
import time
from pathlib import Path
import threading

class ProgressWindow:
    def __init__(self):
        self.window = tk.Toplevel()
        self.window.title("Postęp wyszukiwania")
        self.window.geometry("800x600")
        
        # Ramka dla statystyk
        stats_frame = ttk.LabelFrame(self.window, text="Statystyki", padding=10)
        stats_frame.pack(fill='x', padx=10, pady=5)
        
        # Etykiety dla statystyk
        self.analyzed_files_label = ttk.Label(stats_frame, text="Przeanalizowane pliki: 0")
        self.analyzed_files_label.pack(anchor='w')
        
        self.found_matches_label = ttk.Label(stats_frame, text="Znalezione wystąpienia: 0")
        self.found_matches_label.pack(anchor='w')
        
        self.skipped_files_label = ttk.Label(stats_frame, text="Pominięte pliki: 0")
        self.skipped_files_label.pack(anchor='w')
        
        self.errors_label = ttk.Label(stats_frame, text="Błędy: 0")
        self.errors_label.pack(anchor='w')
        
        # Ramka dla aktualnego postępu
        current_frame = ttk.LabelFrame(self.window, text="Aktualnie przetwarzany plik", padding=10)
        current_frame.pack(fill='x', padx=10, pady=5)
        
        self.current_file_label = ttk.Label(current_frame, text="", wraplength=550)
        self.current_file_label.pack(anchor='w')
        
        # Lista ostatnio przeanalizowanych plików
        list_frame = ttk.LabelFrame(self.window, text="Ostatnio przeanalizowane pliki", padding=10)
        list_frame.pack(fill='both', expand=True, padx=10, pady=5)
        
        # Dodanie scrollbara do listy
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side='right', fill='y')
        
        self.file_listbox = tk.Listbox(list_frame, yscrollcommand=scrollbar.set, width=70, height=10)
        self.file_listbox.pack(fill='both', expand=True)
        scrollbar.config(command=self.file_listbox.yview)
        
        # Przycisk anulowania
        self.cancel_button = ttk.Button(self.window, text="Anuluj wyszukiwanie", command=self.cancel_search)
        self.cancel_button.pack(pady=10)
        
        self.cancelled = False
        self.stats = {
            'analyzed': 0,
            'matches': 0,
            'skipped': 0,
            'errors': 0
        }
        
        # Upewnij się, że okno jest na wierzchu
        self.window.lift()
        self.window.transient()
        
    def update_stats(self, analyzed=0, matches=0, skipped=0, errors=0):
        self.stats['analyzed'] += analyzed
        self.stats['matches'] += matches
        self.stats['skipped'] += skipped
        self.stats['errors'] += errors
        
        self.analyzed_files_label.config(text=f"Przeanalizowane pliki: {self.stats['analyzed']}")
        self.found_matches_label.config(text=f"Znalezione wystąpienia: {self.stats['matches']}")
        self.skipped_files_label.config(text=f"Pominięte pliki: {self.stats['skipped']}")
        self.errors_label.config(text=f"Błędy: {self.stats['errors']}")
        
    def update_current_file(self, filename):
        self.current_file_label.config(text=filename)
        self.file_listbox.insert(0, filename)
        if self.file_listbox.size() > 100:  # Zachowaj tylko ostatnie 100 plików
            self.file_listbox.delete(100)
        self.window.update()
        
    def cancel_search(self):
        self.cancelled = True
        self.cancel_button.config(state='disabled')
        
    def close(self):
        self.window.destroy()

class FileSearchApp:
    def __init__(self):
        self.window = tk.Tk()
        self.window.title("Wyszukiwanie w plikach")
        self.window.geometry("400x350")  # Zwiększona wysokość dla nowego checkboxa
        
        # Ukryj główne okno tk
        self.window.withdraw()
        
        # Utwórz nowe okno jako Toplevel
        self.main_window = tk.Toplevel()
        self.main_window.title("Wyszukiwanie w plikach")
        self.main_window.geometry("400x350")
        
        # Utworzenie interfejsu
        self.create_widgets()
        
    def create_widgets(self):
        # Ramka dla wyboru typów plików
        extensions_frame = ttk.LabelFrame(self.main_window, text="Typy plików do przeszukania", padding=10)
        extensions_frame.pack(fill='x', padx=10, pady=5)
        
        # Checkboxy dla różnych typów plików
        self.doc_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(extensions_frame, text="Dokumenty Word (.doc, .docx)", 
                       variable=self.doc_var).pack(anchor='w')
        
        self.xls_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(extensions_frame, text="Arkusze Excel (.xls, .xlsx)", 
                       variable=self.xls_var).pack(anchor='w')
        
        self.ppt_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(extensions_frame, text="Prezentacje PowerPoint (.ppt, .pptx)", 
                       variable=self.ppt_var).pack(anchor='w')
        
        self.pdf_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(extensions_frame, text="Dokumenty PDF (.pdf)", 
                       variable=self.pdf_var).pack(anchor='w')
        
        # Ramka dla opcji wyszukiwania
        options_frame = ttk.LabelFrame(self.main_window, text="Opcje wyszukiwania", padding=10)
        options_frame.pack(fill='x', padx=10, pady=5)
        
        # Checkbox dla wrażliwości na wielkość liter
        self.case_sensitive_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(options_frame, 
                       text="Uwzględniaj wielkość liter", 
                       variable=self.case_sensitive_var).pack(anchor='w')
        
        # Przycisk rozpoczęcia wyszukiwania
        ttk.Button(self.main_window, text="Rozpocznij wyszukiwanie", 
                  command=self.start_search).pack(pady=20)
        
        # Obsługa zamknięcia okna
        self.main_window.protocol("WM_DELETE_WINDOW", self.on_closing)
        
    def on_closing(self):
        """Obsługa zamknięcia okna"""
        self.window.destroy()
        
    def get_selected_extensions(self):
        """Zwraca listę wybranych rozszerzeń plików"""
        extensions = []
        if self.doc_var.get():
            extensions.extend(['.doc', '.docx'])
        if self.xls_var.get():
            extensions.extend(['.xls', '.xlsx'])
        if self.ppt_var.get():
            extensions.extend(['.ppt', '.pptx'])
        if self.pdf_var.get():
            extensions.append('.pdf')
        return extensions
        
    def is_temp_file(self, filename):
        """Sprawdza czy plik jest plikiem tymczasowym"""
        temp_patterns = [
            '~$', '~',  # Standardowe pliki tymczasowe Office
            '._',       # Pliki tymczasowe macOS
            '.~',       # Inne pliki tymczasowe
            'Backup of' # Kopie zapasowe
        ]
        return any(filename.startswith(pattern) for pattern in temp_patterns)

    def normalize_path(self, path):
        """Normalizuje ścieżkę do pliku"""
        return str(Path(path).resolve())
    
    def is_file_accessible(self, file_path):
        """Sprawdza czy plik jest dostępny do odczytu"""
        try:
            with open(file_path, 'rb'):
                return True
        except:
            return False
        
    def read_word_file(self, file_path, search_phrase, case_sensitive):
        """Przeszukuje plik Word lub XML z rozszerzeniem .doc"""
        results = []
        try:
            # Dodatkowa weryfikacja rozmiaru pliku
            if os.path.getsize(file_path) == 0:
                return [{'error': 'Plik jest pusty', 'type': 'error'}]
                
            # Dodatkowa weryfikacja pliku tymczasowego
            filename = os.path.basename(file_path)
            if filename.startswith('._') or filename.startswith('.~'):
                return [{'error': 'Plik tymczasowy lub ukryty', 'type': 'error'}]
                
            try:
                doc = Document(file_path)
                # Standardowe przeszukiwanie dokumentu Word
                for para in doc.paragraphs:
                    if case_sensitive:
                        if search_phrase in para.text:
                            results.append({
                                'text': para.text,
                                'type': 'paragraph'
                            })
                    else:
                        if search_phrase.lower() in para.text.lower():
                            results.append({
                                'text': para.text,
                                'type': 'paragraph'
                            })
            except Exception as e:
                # Jeśli nie udało się otworzyć jako dokument Word, próbujemy jako XML
                if "not a Word file" in str(e) or "content type" in str(e):
                    try:
                        # Próba odczytu jako zwykły plik tekstowy
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                            content = file.read()
                            if case_sensitive:
                                if search_phrase in content:
                                    start_pos = content.find(search_phrase)
                                    start = max(0, start_pos - 100)
                                    end = min(len(content), start_pos + len(search_phrase) + 100)
                                    context = content[start:end]
                                    results.append({
                                        'text': f"...{context}...",
                                        'type': 'xml_content'
                                    })
                            else:
                                if search_phrase.lower() in content.lower():
                                    start_pos = content.lower().find(search_phrase.lower())
                                    start = max(0, start_pos - 100)
                                    end = min(len(content), start_pos + len(search_phrase) + 100)
                                    context = content[start:end]
                                    results.append({
                                        'text': f"...{context}...",
                                        'type': 'xml_content'
                                    })
                    except Exception as xml_error:
                        return [{'error': f'Nie udało się odczytać pliku ani jako Word ani jako XML: {str(xml_error)}', 'type': 'error'}]
                elif "no relationship" in str(e):
                    return [{'error': 'Plik jest uszkodzony lub ma nieprawidłową strukturę', 'type': 'error'}]
                elif "Package not found" in str(e):
                    return [{'error': 'Nie można otworzyć pliku - może być uszkodzony lub zablokowany', 'type': 'error'}]
                else:
                    return [{'error': str(e), 'type': 'error'}]
                    
            return results
            
        except Exception as e:
            return [{'error': f'Nieoczekiwany błąd: {str(e)}', 'type': 'error'}]

    def read_excel_file(self, file_path, search_phrase, case_sensitive):
        """Przeszukuje plik Excel"""
        results = []
        try:
            # Próba odczytu pliku Excel przy użyciu pandas
            try:
                df = pd.read_excel(file_path)
            except ImportError as e:
                if "xlrd" in str(e):
                    # Jeśli brak xlrd, próbujemy odczytać plik jako CSV
                    try:
                        with open(file_path, 'rb') as file:
                            content = file.read().decode('utf-8', errors='ignore')
                            if case_sensitive:
                                if search_phrase in content:
                                    start_pos = content.find(search_phrase)
                                    start = max(0, start_pos - 100)
                                    end = min(len(content), start_pos + len(search_phrase) + 100)
                                    context = content[start:end]
                                    results.append({
                                        'text': f"...{context}...",
                                        'type': 'raw_content'
                                    })
                            else:
                                if search_phrase.lower() in content.lower():
                                    start_pos = content.lower().find(search_phrase.lower())
                                    start = max(0, start_pos - 100)
                                    end = min(len(content), start_pos + len(search_phrase) + 100)
                                    context = content[start:end]
                                    results.append({
                                        'text': f"...{context}...",
                                        'type': 'raw_content'
                                    })
                            return results
                    except Exception as read_error:
                        return [{'error': f'Nie można odczytać pliku Excel: {str(read_error)}', 'type': 'error'}]
                else:
                    raise e

            # Standardowe przeszukiwanie jeśli udało się odczytać jako Excel
            for col in df.columns:
                matches = df[df[col].astype(str).str.contains(
                    search_phrase, 
                    case=case_sensitive, 
                    na=False, 
                    regex=False
                )]
                if not matches.empty:
                    for _, row in matches.iterrows():
                        results.append({
                            'text': f"Kolumna: {col}, Wartość: {row[col]}",
                            'type': 'cell'
                        })
            return results
            
        except Exception as e:
            return [{'error': str(e), 'type': 'error'}]

    def read_powerpoint_file(self, file_path, search_phrase, case_sensitive):
        """Przeszukuje plik PowerPoint"""
        results = []
        try:
            try:
                prs = Presentation(file_path)
            except Exception as ppt_error:
                # Jeśli nie udało się otworzyć jako PPTX, próbujemy odczytać zawartość binarnie
                try:
                    with open(file_path, 'rb') as file:
                        content = file.read().decode('utf-8', errors='ignore')
                        if case_sensitive:
                            if search_phrase in content:
                                start_pos = content.find(search_phrase)
                                start = max(0, start_pos - 100)
                                end = min(len(content), start_pos + len(search_phrase) + 100)
                                context = content[start:end]
                                results.append({
                                    'text': f"...{context}...",
                                    'type': 'raw_content'
                                })
                        else:
                            if search_phrase.lower() in content.lower():
                                start_pos = content.lower().find(search_phrase.lower())
                                start = max(0, start_pos - 100)
                                end = min(len(content), start_pos + len(search_phrase) + 100)
                                context = content[start:end]
                                results.append({
                                    'text': f"...{context}...",
                                    'type': 'raw_content'
                                })
                        return results
                except Exception as read_error:
                    return [{'error': f'Nie można odczytać pliku PowerPoint: {str(read_error)}', 'type': 'error'}]
                return [{'error': str(ppt_error), 'type': 'error'}]

            # Standardowe przeszukiwanie jeśli udało się otworzyć jako PPTX
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if case_sensitive:
                            if search_phrase in shape.text:
                                results.append({
                                    'text': shape.text,
                                    'type': 'slide'
                                })
                        else:
                            if search_phrase.lower() in shape.text.lower():
                                results.append({
                                    'text': shape.text,
                                    'type': 'slide'
                                })
            return results
            
        except Exception as e:
            return [{'error': str(e), 'type': 'error'}]

    def read_pdf_file(self, file_path, search_phrase, case_sensitive):
        """Przeszukuje plik PDF"""
        results = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    text = pdf_reader.pages[page_num].extract_text()
                    if case_sensitive:
                        if search_phrase in text:
                            results.append({
                                'text': f"Strona {page_num + 1}: {text[:200]}...",
                                'type': 'page'
                            })
                    else:
                        if search_phrase.lower() in text.lower():
                            results.append({
                                'text': f"Strona {page_num + 1}: {text[:200]}...",
                                'type': 'page'
                            })
            return results
        except Exception as e:
            return [{'error': str(e), 'type': 'error'}]

    def search_files(self, directory, search_phrase, extensions):
        """Przeszukuje pliki we wskazanym katalogu"""
        results = {}
        skipped_files = []
        inaccessible_files = []
        
        # Pobierz stan wrażliwości na wielkość liter
        case_sensitive = self.case_sensitive_var.get()
        
        # Utworzenie okna postępu
        progress_window = ProgressWindow()
        
        try:
            for root, _, files in os.walk(directory):
                for file in files:
                    if progress_window.cancelled:
                        messagebox.showinfo("Anulowano", "Wyszukiwanie zostało anulowane przez użytkownika.")
                        progress_window.close()
                        return results, skipped_files, inaccessible_files
                    
                    try:
                        if self.is_temp_file(file):
                            skipped_files.append(os.path.join(root, file))
                            progress_window.update_stats(skipped=1)
                            continue
                            
                        file_path = self.normalize_path(os.path.join(root, file))
                        file_ext = os.path.splitext(file)[1].lower()
                        
                        if file_ext not in extensions:
                            continue
                        
                        # Aktualizacja okna postępu
                        progress_window.update_current_file(file_path)
                        
                        if not self.is_file_accessible(file_path):
                            inaccessible_files.append(file_path)
                            progress_window.update_stats(errors=1)
                            continue
                        
                        # Przekazanie informacji o wrażliwości na wielkość liter do funkcji przeszukujących
                        file_results = []
                        if file_ext in ['.doc', '.docx']:
                            file_results = self.read_word_file(file_path, search_phrase, case_sensitive)
                        elif file_ext in ['.xls', '.xlsx']:
                            file_results = self.read_excel_file(file_path, search_phrase, case_sensitive)
                        elif file_ext in ['.ppt', '.pptx']:
                            file_results = self.read_powerpoint_file(file_path, search_phrase, case_sensitive)
                        elif file_ext == '.pdf':
                            file_results = self.read_pdf_file(file_path, search_phrase, case_sensitive)
                        
                        # Sprawdzenie czy są błędy w wynikach
                        has_errors = any('error' in r for r in file_results)
                        if has_errors:
                            progress_window.update_stats(errors=1)
                            
                        if file_results:
                            results[file_path] = file_results
                            matches_count = sum(1 for r in file_results if 'error' not in r)
                            if matches_count > 0:  # Aktualizuj statystyki tylko jeśli znaleziono wystąpienia
                                progress_window.update_stats(analyzed=1, matches=matches_count)
                            else:
                                progress_window.update_stats(analyzed=1)
                        else:
                            progress_window.update_stats(analyzed=1)
                                
                    except Exception as e:
                        inaccessible_files.append(f"{file_path} (Błąd: {str(e)})")
                        progress_window.update_stats(errors=1)
                        
        except Exception as e:
            messagebox.showerror("Błąd", f"Wystąpił błąd podczas przeszukiwania katalogu: {str(e)}")
        
        progress_window.close()
        return results, skipped_files, inaccessible_files

    def save_results(self, results, skipped_files, inaccessible_files, search_phrase):
        """Zapisuje wyniki do pliku"""
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = f"wyniki_wyszukiwania_{timestamp}.txt"
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"Wyniki wyszukiwania frazy: {search_phrase}\n")
            f.write(f"Data wyszukiwania: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            
            # 1. Najpierw wyświetlamy znalezione wystąpienia (bez błędów)
            successful_results = {}
            error_results = {}
            
            for file_path, file_results in results.items():
                # Rozdzielamy wyniki na udane i błędy
                has_success = any('error' not in r for r in file_results)
                if has_success:
                    successful_results[file_path] = [r for r in file_results if 'error' not in r]
                if any('error' in r for r in file_results):
                    error_results[file_path] = [r for r in file_results if 'error' in r]
            
            if successful_results:
                f.write("ZNALEZIONE WYSTĄPIENIA:\n")
                f.write("=" * 50 + "\n\n")
                for file_path, file_results in successful_results.items():
                    f.write(f"Plik: {file_path}\n")
                    for result in file_results:
                        if result['type'] == 'xml_content':
                            f.write("Typ: Zawartość pliku XML/DOC\n")
                        else:
                            f.write(f"Typ: {result['type']}\n")
                        f.write(f"Tekst: {result['text']}\n")
                    f.write("-" * 50 + "\n\n")
            
            # 2. Następnie wyświetlamy pominięte pliki
            if skipped_files:
                f.write("\nPOMINIĘTE PLIKI TYMCZASOWE:\n")
                f.write("=" * 50 + "\n")
                for file in skipped_files:
                    f.write(f"{file}\n")
                f.write("\n")
            
            # 3. Podsumowanie
            f.write("\nPODSUMOWANIE:\n")
            f.write("=" * 50 + "\n")
            f.write(f"Liczba plików z znalezionymi wystąpieniami: {len(successful_results)}\n")
            f.write(f"Liczba pominiętych plików tymczasowych: {len(skipped_files)}\n")
            f.write(f"Liczba plików z błędami: {len(error_results) + len(inaccessible_files)}\n\n")
            
            # 4. Informacja o xlrd
            f.write("\nINFORMACJA:\n")
            f.write("Aby poprawić obsługę starszych plików Excel (.xls), zainstaluj bibliotekę xlrd:\n")
            f.write("pip install xlrd>=2.0.1\n\n")
            
            # 5. Na końcu wyświetlamy wszystkie błędy
            if error_results or inaccessible_files:
                f.write("\nPLIKI Z BŁĘDAMI:\n")
                f.write("=" * 50 + "\n")
                
                # Błędy z przetwarzania plików
                for file_path, file_results in error_results.items():
                    f.write(f"Plik: {file_path}\n")
                    for result in file_results:
                        if 'error' in result:
                            f.write(f"BŁĄD: {result['error']}\n")
                    f.write("-" * 50 + "\n")
                
                # Błędy dostępu do plików
                if inaccessible_files:
                    for file in inaccessible_files:
                        f.write(f"{file}\n")
                        f.write("-" * 50 + "\n")
                        
        return output_file

    def start_search(self):
        """Rozpoczyna proces wyszukiwania"""
        self.main_window.destroy()
        
        directory = filedialog.askdirectory(title="Wybierz katalog do przeszukania")
        if not directory:
            messagebox.showerror("Błąd", "Nie wybrano katalogu!")
            self.window.destroy()  # Zamknij aplikację
            return
            
        search_phrase = tk.simpledialog.askstring("Wyszukiwanie", "Wprowadź frazę do wyszukania:")
        if not search_phrase:
            messagebox.showerror("Błąd", "Nie wprowadzono frazy do wyszukania!")
            self.window.destroy()  # Zamknij aplikację
            return
            
        extensions = self.get_selected_extensions()
        if not extensions:
            messagebox.showerror("Błąd", "Nie wybrano żadnych typów plików!")
            self.window.destroy()  # Zamknij aplikację
            return
        
        results, skipped_files, inaccessible_files = self.search_files(directory, search_phrase, extensions)
        
        output_file = self.save_results(results, skipped_files, inaccessible_files, search_phrase)
        
        # Poprawione zliczanie wystąpień i błędów
        found_occurrences = 0
        files_with_phrase = 0
        files_with_errors = len(inaccessible_files)  # Pliki niedostępne
        
        for file_path, file_results in results.items():
            # Sprawdzamy czy plik zawiera błędy
            if any('error' in r for r in file_results):
                files_with_errors += 1
            else:
                # Zliczamy tylko wyniki bez błędów
                valid_results = [r for r in file_results if 'error' not in r]
                if valid_results:
                    files_with_phrase += 1
                    found_occurrences += len(valid_results)
        
        message = f"Wyszukiwanie zakończone!\n\n"
        message += f"Przeanalizowano pliki z rozszerzeniami: {', '.join(extensions)}\n"
        
        if found_occurrences > 0:
            message += f"Znaleziono {found_occurrences} wystąpień frazy w {files_with_phrase} plikach.\n"
        else:
            message += "Nie znaleziono żadnych wystąpień szukanej frazy.\n"
            
        message += f"Pominięto {len(skipped_files)} plików tymczasowych.\n"
        message += f"Nie udało się przeanalizować {files_with_errors} plików.\n\n"
        message += f"Wyniki zapisano do pliku:\n{output_file}"
        
        # Wyświetl komunikat i zakończ aplikację
        def on_close_info():
            self.window.quit()
            self.window.destroy()
        
        info_window = tk.Toplevel()
        info_window.withdraw()  # Ukryj puste okno
        info_window.after(1, lambda: [
            messagebox.showinfo("Zakończono", message),
            on_close_info()
        ])

if __name__ == "__main__":
    import sys
    
    # Przekierowanie stderr do null (ukrycie komunikatów Google Drive)
    if sys.platform == 'win32':
        sys.stderr = open('nul', 'w')
    else:
        sys.stderr = open('/dev/null', 'w')
    
    try:
        app = FileSearchApp()
        app.window.mainloop()
    except Exception as e:
        messagebox.showerror("Błąd", f"Wystąpił nieoczekiwany błąd: {str(e)}")
    finally:
        # Upewnij się, że aplikacja zostanie zamknięta
        try:
            app.window.destroy()
        except:
            pass
        sys.exit(0)
